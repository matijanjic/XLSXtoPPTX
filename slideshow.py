from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import PIL
import io


class SlideShow:

    def __init__(self, width, height, layoutNumber):
        
        self.height = height
        self.width = width
        self.layoutNumber = layoutNumber
        self.ss = Presentation()
        self.slideLayout = self.ss.slide_layouts[layoutNumber]
        self.ss.slide_width = Inches(self.width)
        self.ss.slide_height = Inches(self.height)

    # adds an empty slide
    def addSlide(self):
        self.slide = self.ss.slides.add_slide(self.slideLayout)

    # adds a text object on the slide
    def addText(self, fontSize, text, width, height, left='center', top='center'):
        if not isinstance(text, str):
            print("addText function requires a string, entered value is " + str(type(text)))
            exit(1)
        # convert width and height to EMUs
        width = Inches(width)
        height = Inches(height)
        

        # if left and top kwargs not 'center', convert them to EMUs
        if not left == 'center':
            left = Inches(left)
        if not top == 'center':
            top = Inches(top)

        # if kwargs left and top equal to 'center', text is centered so the user doesn't have to convert when calling the method
        if left == 'center':
            left = self.ss.slide_width / 2
            left = left - (width / 2)
        if top == 'center':
            top = self.ss.slide_height / 2
            top = top - (height / 2)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.text = text
        p.font.size = Pt(fontSize)

    # add a picture on the slide, maxSize is defined in pixels. Default is centered, but if needed, the picture
    # can be moved anywhere, just need to specify the left and/or top variables in inches
    def addPicture(self, imgFile, maxSize, left='center', top='center'):
        # if left and/or top kwargs not 'center', convert them to EMUs
        if not left == 'center':
            left = Inches(left)
        if not top == 'center':
            top = Inches(top)

        # open the image
        img = Image.open(imgFile)

        # depending on the picture orientation, set the longest side to the maxSize argument (in pixels)
        width, height = img.size

        if height > width:
            ratio = width / height
            height = maxSize
            width = int(maxSize * ratio)
        else:
            ratio = height / width
            width = maxSize
            height = int(maxSize * ratio)

        # resize the image using the calculations above and save it to img_resized variable
        img_resized = img.resize([width, height], PIL.Image.ANTIALIAS)

        # using BytesIO to save the resized image to memory
        with io.BytesIO() as output:
            img_resized.save(output, img.format)

            # picture can be centered vertically, horizontally or both if so selected with the kwargs
            if left == 'center' and top == 'center':
                pic = self.slide.shapes.add_picture(output, 0, 0)
                pic.left = int((self.ss.slide_width - pic.width)/2)
                pic.top = int((self.ss.slide_height - pic.height)/2)
            elif left == 'center' and not top == 'center':
                pic = self.slide.shapes.add_picture(output, 0, 0)
                pic.left = int((self.ss.slide_width - pic.width)/2)
                pic.top = top
            elif not left == 'center' and top == 'center':
                pic = self.slide.shapes.add_picture(output, 0, 0)
                pic.top = int((self.ss.slide_height - pic.height)/2)
                pic.left = left
            # else position it depending on the left and top variable (in inches)
            else:
                self.slide.shapes.add_picture(output, left, top)

    def addSound(self, soundFile):
        # uses the add_movie method of the pptx module
        sound = self.slide.shapes.add_movie(soundFile, 0, 0, 0, 0)
        
        # a bit of xml editing using the etree method from the lxml module making sound autoplay possible.
        # Solution found at https://github.com/scanny/python-pptx/issues/427. Thanks to iota-pi for the solution!
        tree = sound._element.getparent().getparent().getnext().getnext()
        timing = [el for el in tree.iterdescendants(
        ) if etree.QName(el).localname == 'cond'][0]
        timing.set('delay', '0')

    # saves the pptx file
    def save(self, saveFile):
        self.ss.save(saveFile)
